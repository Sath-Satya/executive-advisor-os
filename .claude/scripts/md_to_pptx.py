"""
md_to_pptx.py — Convert Marp-style Markdown deck to cream-navy-gold PPTX.

Usage:
    python md_to_pptx.py <input.md> [<output.pptx>]
    python md_to_pptx.py <input.md> --auto-install    # install python-pptx if missing

If `python-pptx` is not installed, the script prints the exact install command
and exits 2. Pass `--auto-install` or set `EAOS_AUTO_INSTALL=1` to auto-install.

Rules for input:
- Slides separated by horizontal rule `---` on its own line (after the YAML frontmatter)
- Marp YAML frontmatter at top of file is stripped and ignored (python-pptx applies the theme)
- Presenter notes: content inside `<!-- ... -->` at end of a slide becomes speaker notes
- `<!-- _class: lead -->` marks a hero / title-style slide (navy background, cream text)

Palette (per Executive Advisor OS visual-identity rule):
    cream  #f7f4ef      navy   #0e1b2e      gold   #b8965a
    cream2 #ede9e2      navy2  #152440      gold2  #d4ab72
    ink    #1c1a18      muted  #6b6560

Fonts:
    Headings — Cormorant Garamond (serif) with Georgia fallback
    Body     — DM Sans (sans) with system fallback
"""

import re
import sys
import subprocess
from pathlib import Path


def _ensure_pptx():
    """Import python-pptx. If missing, print a clear install hint and offer
    to auto-install with user consent (env var or --auto-install flag).
    """
    try:
        import pptx  # noqa: F401
        return
    except ImportError:
        pass

    msg = (
        "\n"
        "──────────────────────────────────────────────────────────────\n"
        "  Missing dependency: python-pptx\n"
        "──────────────────────────────────────────────────────────────\n"
        "  This converter needs the `python-pptx` package to build\n"
        "  PowerPoint files from the Marp Markdown source.\n"
        "\n"
        "  Install once:\n"
        "      pip install python-pptx\n"
        "\n"
        "  Or install everything the repo needs:\n"
        "      pip install -r requirements.txt\n"
        "──────────────────────────────────────────────────────────────\n"
    )

    auto = "--auto-install" in sys.argv or __import__("os").environ.get("EAOS_AUTO_INSTALL") == "1"
    if auto:
        print(msg, flush=True)
        print("[auto-install] Running: pip install python-pptx", flush=True)
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        except subprocess.CalledProcessError as e:
            print(f"[auto-install failed] {e}", flush=True)
            sys.exit(2)
        try:
            import pptx  # noqa: F401
            if "--auto-install" in sys.argv:
                sys.argv.remove("--auto-install")
            return
        except ImportError:
            print("[auto-install] install completed but import still fails", flush=True)
            sys.exit(2)
    else:
        print(msg, flush=True)
        sys.exit(2)


_ensure_pptx()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------- Palette ----------
CREAM = RGBColor(0xF7, 0xF4, 0xEF)
CREAM2 = RGBColor(0xED, 0xE9, 0xE2)
NAVY = RGBColor(0x0E, 0x1B, 0x2E)
NAVY2 = RGBColor(0x15, 0x24, 0x40)
GOLD = RGBColor(0xB8, 0x96, 0x5A)
GOLD2 = RGBColor(0xD4, 0xAB, 0x72)
INK = RGBColor(0x1C, 0x1A, 0x18)
MUTED = RGBColor(0x6B, 0x65, 0x60)

FONT_HEAD = "Cormorant Garamond"
FONT_HEAD_FALLBACK = "Georgia"
FONT_BODY = "DM Sans"
FONT_BODY_FALLBACK = "Segoe UI"


# ---------- Markdown parsing ----------
def strip_frontmatter(text: str) -> str:
    """Remove Marp YAML frontmatter (--- ... ---) from the top of the file."""
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            # Skip past the closing ---\n
            newline = text.find("\n", end + 4)
            return text[newline + 1 :] if newline != -1 else ""
    return text


def split_slides(text: str):
    """Split markdown into slides using `---` separators at line start.

    Returns a list of slide text blocks.
    """
    # Split on lines consisting entirely of "---" (horizontal rule in Markdown / Marp slide separator)
    parts = re.split(r"(?:\r?\n)---\s*(?:\r?\n)", "\n" + text)
    # Remove leading newline artifact on first part
    slides = [s.strip() for s in parts if s.strip()]
    return slides


def extract_notes(slide_text: str):
    """Extract presenter notes from trailing `<!-- ... -->` comment blocks.

    Returns (cleaned_slide_text, notes_string_or_empty).
    """
    notes_parts = []
    cleaned = slide_text
    # Repeatedly pull the last HTML comment out if it's at end of slide (standard Marp convention)
    comment_re = re.compile(r"<!--\s*(.*?)\s*-->", flags=re.DOTALL)
    for m in comment_re.finditer(slide_text):
        inner = m.group(1).strip()
        # If this comment contains "_class:" it's a Marp slide directive, not notes
        if inner.lower().startswith("_class"):
            continue
        notes_parts.append(inner)
    cleaned = comment_re.sub("", cleaned).strip()
    return cleaned, "\n\n".join(notes_parts).strip()


def is_lead_slide(slide_text: str) -> bool:
    """Check for Marp `<!-- _class: lead -->` directive."""
    return bool(re.search(r"<!--\s*_class:\s*lead\s*-->", slide_text))


def parse_slide_blocks(slide_text: str):
    """Break a slide into structured blocks — h1, h2, h3, paragraphs, tables, lists, blockquotes.

    Returns a list of (kind, content) tuples.
    kind ∈ {"h1", "h2", "h3", "h4", "p", "ul", "ol", "quote", "table", "hr"}
    """
    blocks = []
    lines = slide_text.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue
        # Headings
        if stripped.startswith("#### "):
            blocks.append(("h4", stripped[5:].strip()))
            i += 1
        elif stripped.startswith("### "):
            blocks.append(("h3", stripped[4:].strip()))
            i += 1
        elif stripped.startswith("## "):
            blocks.append(("h2", stripped[3:].strip()))
            i += 1
        elif stripped.startswith("# "):
            blocks.append(("h1", stripped[2:].strip()))
            i += 1
        elif stripped.startswith("> "):
            # Blockquote — collect contiguous
            quote_lines = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote_lines.append(lines[i].strip().lstrip(">").strip())
                i += 1
            blocks.append(("quote", "\n".join(quote_lines).strip()))
        elif stripped.startswith("|") and i + 1 < len(lines) and re.match(r"^\|[\s:|-]+\|$", lines[i + 1].strip()):
            # Table — collect contiguous pipe-rows
            table_lines = [lines[i]]
            i += 1
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            blocks.append(("table", "\n".join(table_lines)))
        elif re.match(r"^\s*[-*]\s+", line):
            # Unordered list — collect contiguous
            list_lines = []
            while i < len(lines) and (re.match(r"^\s*[-*]\s+", lines[i]) or (lines[i].startswith("  ") and lines[i].strip())):
                m = re.match(r"^\s*[-*]\s+(.*)", lines[i])
                if m:
                    list_lines.append(m.group(1).strip())
                i += 1
            blocks.append(("ul", list_lines))
        elif re.match(r"^\s*\d+\.\s+", line):
            list_lines = []
            while i < len(lines) and re.match(r"^\s*\d+\.\s+", lines[i]):
                m = re.match(r"^\s*\d+\.\s+(.*)", lines[i])
                if m:
                    list_lines.append(m.group(1).strip())
                i += 1
            blocks.append(("ol", list_lines))
        else:
            # Paragraph — consume lines until we hit a block-starting pattern
            # A block-starting pattern is a heading, blockquote, pipe-table, bullet list, or numbered list.
            # Leading asterisks/dashes inside a paragraph (e.g. **bold**) do NOT start a block.
            para_lines = [lines[i].strip()]
            i += 1
            while i < len(lines):
                nxt = lines[i]
                s = nxt.strip()
                if not s:
                    break
                if re.match(r"^#{1,6}\s", s):
                    break
                if s.startswith("> "):
                    break
                if s.startswith("|") and i + 1 < len(lines) and re.match(r"^\|[\s:|-]+\|$", lines[i + 1].strip()):
                    break
                if re.match(r"^\s*[-*]\s+", nxt):
                    break
                if re.match(r"^\s*\d+\.\s+", nxt):
                    break
                para_lines.append(s)
                i += 1
            blocks.append(("p", " ".join(para_lines)))
    return blocks


def parse_table_block(raw: str):
    rows = [r for r in raw.split("\n") if r.strip().startswith("|")]
    parsed = []
    for r in rows:
        cells = [c.strip() for c in r.strip().strip("|").split("|")]
        parsed.append(cells)
    # Drop the separator row (all dashes/colons)
    return [r for r in parsed if not all(re.match(r"^[\s:-]+$", c) for c in r)]


# ---------- Inline markdown → runs ----------
INLINE_PATTERN = re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*|__[^_]+__|_[^_]+_|`[^`]+`)")


def add_inline(paragraph, text: str, base_color=INK, font=FONT_BODY, size=18):
    """Add runs with inline markdown — **bold**, *italic*, `code`."""
    if not text:
        return
    pos = 0
    for m in INLINE_PATTERN.finditer(text):
        if m.start() > pos:
            run = paragraph.add_run()
            run.text = text[pos : m.start()]
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = base_color
        token = m.group(0)
        run = paragraph.add_run()
        if token.startswith("**") or token.startswith("__"):
            run.text = token[2:-2]
            run.font.bold = True
            run.font.color.rgb = NAVY if base_color == INK else base_color
        elif token.startswith("`"):
            run.text = token[1:-1]
            run.font.name = "Consolas"
            run.font.color.rgb = GOLD
        else:
            run.text = token[1:-1]
            run.font.italic = True
            run.font.color.rgb = GOLD if base_color == INK else base_color
        run.font.size = Pt(size)
        if not run.font.name:
            run.font.name = font
        pos = m.end()
    if pos < len(text):
        run = paragraph.add_run()
        run.text = text[pos:]
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = base_color


# ---------- Slide builders ----------
def set_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer_bar(slide, prs, is_lead=False):
    """Thin gold accent bar at top of slide."""
    bar_color = GOLD
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Emu(32000)  # ~0.03 inches
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()


def build_lead_slide(prs, blocks, notes, slide_num):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    set_background(slide, NAVY)
    add_footer_bar(slide, prs)

    # Build vertically-centered title + subtitle + bullets
    left = Inches(0.8)
    top = Inches(0.9)
    width = prs.slide_width - Inches(1.6)
    height = prs.slide_height - Inches(1.8)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for kind, content in blocks:
        if kind == "h1":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = content
            run.font.name = FONT_HEAD
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = CREAM
            p.space_after = Pt(14)
        elif kind == "h2":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = content
            run.font.name = FONT_HEAD
            run.font.size = Pt(26)
            run.font.italic = True
            run.font.color.rgb = GOLD2
            p.space_after = Pt(18)
        elif kind == "h3":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = content.upper()
            run.font.name = FONT_BODY
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = GOLD
            p.space_after = Pt(10)
        elif kind == "p":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            add_inline(p, content, base_color=CREAM, font=FONT_BODY, size=16)
            p.space_after = Pt(8)
        elif kind == "ul":
            for item in content:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                # Add gold bullet char manually since python-pptx bullet formatting is fiddly
                add_inline(p, "•  " + item, base_color=CREAM, font=FONT_BODY, size=15)
                p.space_after = Pt(4)
    # Slide number in bottom right
    _add_slide_number(slide, prs, slide_num, on_dark=True)
    set_notes(slide, notes)
    return slide


def build_content_slide(prs, blocks, notes, slide_num):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_background(slide, CREAM)
    add_footer_bar(slide, prs)

    left = Inches(0.55)
    top = Inches(0.4)
    width = prs.slide_width - Inches(1.1)
    height = prs.slide_height - Inches(0.8)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    first = True

    def add_para():
        nonlocal first
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        return p

    for kind, content in blocks:
        if kind == "h1":
            p = add_para()
            run = p.add_run()
            run.text = content
            run.font.name = FONT_HEAD
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = NAVY
            p.space_after = Pt(10)
        elif kind == "h2":
            p = add_para()
            run = p.add_run()
            run.text = content
            run.font.name = FONT_HEAD
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = NAVY
            p.space_after = Pt(8)
        elif kind == "h3":
            p = add_para()
            run = p.add_run()
            run.text = content.upper()
            run.font.name = FONT_BODY
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = GOLD
            p.space_after = Pt(6)
        elif kind == "h4":
            p = add_para()
            run = p.add_run()
            run.text = content
            run.font.name = FONT_BODY
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = NAVY
            p.space_after = Pt(4)
        elif kind == "p":
            p = add_para()
            add_inline(p, content, base_color=INK, font=FONT_BODY, size=14)
            p.space_after = Pt(6)
        elif kind == "ul":
            for item in content:
                p = add_para()
                add_inline(p, "•  " + item, base_color=INK, font=FONT_BODY, size=13)
                p.space_after = Pt(3)
        elif kind == "ol":
            for idx, item in enumerate(content, 1):
                p = add_para()
                add_inline(p, f"{idx}.  " + item, base_color=INK, font=FONT_BODY, size=13)
                p.space_after = Pt(3)
        elif kind == "quote":
            p = add_para()
            add_inline(p, content, base_color=NAVY, font=FONT_HEAD, size=17)
            p.space_after = Pt(8)
        elif kind == "table":
            rows = parse_table_block(content)
            if rows:
                _add_table(slide, rows, left, Inches(0.3 + 0.1 * len(blocks)), prs.slide_width)

    _add_slide_number(slide, prs, slide_num, on_dark=False)
    set_notes(slide, notes)
    return slide


def _add_table(slide, rows, left, top, slide_width_emu):
    ncols = max(len(r) for r in rows)
    nrows = len(rows)
    width = slide_width_emu - Inches(1.3)
    height = Emu(int(300000 * nrows))
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            if c_idx >= ncols:
                continue
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            add_inline(p, val, base_color=NAVY if r_idx == 0 else INK, font=FONT_BODY, size=10)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM2
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_slide_number(slide, prs, num, on_dark):
    box_w = Inches(0.8)
    box_h = Inches(0.3)
    left = prs.slide_width - box_w - Inches(0.3)
    top = prs.slide_height - box_h - Inches(0.25)
    tb = slide.shapes.add_textbox(left, top, box_w, box_h)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num:02d}"
    run.font.name = FONT_BODY
    run.font.size = Pt(9)
    run.font.color.rgb = GOLD2 if on_dark else MUTED


def set_notes(slide, notes):
    if not notes:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes


# ---------- Main ----------
def convert(md_path: Path, pptx_path: Path):
    text = md_path.read_text(encoding="utf-8")
    text = strip_frontmatter(text)
    raw_slides = split_slides(text)
    if not raw_slides:
        raise ValueError(f"No slides found in {md_path}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, raw in enumerate(raw_slides, 1):
        cleaned, notes = extract_notes(raw)
        blocks = parse_slide_blocks(cleaned)
        if is_lead_slide(raw) or idx == 1:
            build_lead_slide(prs, blocks, notes, idx)
        else:
            build_content_slide(prs, blocks, notes, idx)

    prs.save(str(pptx_path))
    print(f"[OK] {pptx_path}  ({len(raw_slides)} slides)")


def main():
    args = [a for a in sys.argv[1:] if a != "--auto-install"]
    if len(args) < 1:
        print("Usage: python md_to_pptx.py <input.md> [<output.pptx>] [--auto-install]")
        sys.exit(1)
    md_path = Path(args[0]).resolve()
    if not md_path.exists():
        print(f"[ERROR] {md_path} not found")
        sys.exit(1)
    if len(args) >= 2:
        pptx_path = Path(args[1]).resolve()
    else:
        pptx_path = md_path.with_suffix(".pptx")
    convert(md_path, pptx_path)


if __name__ == "__main__":
    main()
